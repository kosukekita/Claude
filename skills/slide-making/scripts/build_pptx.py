# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx"]
# ///
"""
Build a native .pptx deck from a JSON spec — WITHOUT going through HTML.

This is the PPTX path of the slide-making skill. Claude's job is to author the
JSON spec from the user's draft; this script owns all the mechanical python-pptx
work (EMU math, run-level formatting, the design system, horizontal-only tables,
placing GPT-Image part PNGs). Mirrors render_slide.py's split: the model writes
content, the script renders it deterministically.

Usage:
    uv run --with python-pptx build_pptx.py --spec deck.json --out deck.pptx
    cat deck.json | uv run --with python-pptx build_pptx.py --spec - --out deck.pptx

Spec schema (JSON):
{
  "slides": [
    {
      "title": "現状の課題",                 // required
      "subtitle": "2025 Q2 レビュー",         // optional, under the title
      "bullets": ["コスト上昇", "解約率増"],   // optional, "・" prepended
      "body": "段落テキスト",                 // optional, plain paragraph
      "page_number": 3,                       // optional, bottom-right
      "images": [                             // optional GPT-Image/theSVG parts
        {"path": "parts/icon.png", "x_in": 1.0, "y_in": 3.2, "w_in": 1.2}
        // omit h_in to preserve aspect ratio
      ],
      "table": {                              // optional, horizontal lines only
        "headers": ["指標", "前年", "今年"],
        "rows": [["CAC", "¥8,200", "¥11,500"]],
        "x_in": 7.0, "y_in": 2.5, "w_in": 5.5, "h_in": 2.5,
        "emphasize_row": 1                    // optional 1-based data row index
      },
      "emphasis": [                           // optional, recolors matching runs
        {"text": "1.4倍", "kind": "accent"}   // kind: main|accent|underline|invert
      ]
    }
  ]
}

Image paths in the spec are resolved relative to the spec file's directory
(or CWD when reading from stdin), so a deck folder is portable across machines.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# --- Design system (mirrors references/design-rules.md) ----------------------
BASE = RGBColor(0xF9, 0xF9, 0xF9)    # slide background, ~70%
TEXT = RGBColor(0x1A, 0x1A, 0x1A)    # headings/body/table rules, ~25%
MAIN = RGBColor(0x00, 0x71, 0xBC)    # main accent, keep <4%
ACCENT = RGBColor(0xFF, 0x50, 0x50)  # warning/negative, last resort <1%

FONT = "Noto Sans JP"                # no font-embedding in pptx; opener substitutes
TITLE_PT, HEAD_PT, BODY_PT, NUM_PT = Pt(50), Pt(35), Pt(25), Pt(18)

# 1920x1080 @96dpi -> 13.333in x 7.5in; 48px margin -> 0.5in
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.5)
CONTENT_W = SLIDE_W - 2 * MARGIN

EMPHASIS = {"main": MAIN, "accent": ACCENT}  # underline/invert handled separately


def _set_run(run, *, size=BODY_PT, bold=False, color=TEXT, underline=False):
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    run.font.underline = underline
    run.font.color.rgb = color
    # Ensure CJK (east-asian) font is also FONT, else PowerPoint may swap it.
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def _estimate_lines(text, size_pt, box_w_in):
    """Rough estimate of how many lines `text` wraps to in a box `box_w_in` wide.

    Used only to push following content down so a long subtitle doesn't overlap
    the bullets. CJK chars count as ~1 em, ASCII as ~0.55 em. Em width ~= size in
    points / 72 inches. Deliberately conservative (rounds up).
    """
    pt = size_pt.pt if hasattr(size_pt, "pt") else float(size_pt)
    em_in = pt / 72.0
    units = 0.0
    for ch in str(text):
        units += 0.55 if ord(ch) < 0x3000 else 1.0  # ascii/punct vs CJK/full-width
    text_w_in = units * em_in
    import math
    return max(1, math.ceil(text_w_in / max(box_w_in, 0.1)))


def _add_textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _emphasis_lookup(text, emphasis):
    """Return the first emphasis spec whose token equals `text`, else None."""
    for emp in emphasis:
        if emp.get("text", "") == text:
            return emp
    return None


def _split_on_tokens(text, tokens):
    """Split `text` into segments, keeping emphasis tokens as their own pieces.

    Longest tokens first so overlapping tokens resolve to the longer match.
    Returns a list of substrings; a substring equal to a token is an emphasis
    segment, everything else is plain. Adjacent plain text stays merged.
    """
    if not tokens:
        return [text]
    pieces = [text]
    for tok in sorted(tokens, key=len, reverse=True):
        if not tok:
            continue
        nxt = []
        for piece in pieces:
            if piece in tokens:  # already an emphasis segment, don't re-split
                nxt.append(piece)
                continue
            start = 0
            while True:
                idx = piece.find(tok, start)
                if idx == -1:
                    nxt.append(piece[start:])
                    break
                if idx > start:
                    nxt.append(piece[start:idx])
                nxt.append(tok)
                start = idx + len(tok)
        pieces = [p for p in nxt if p != ""]
    return pieces


def _apply_emphasis(tf, emphasis):
    """Recolor ONLY the emphasized tokens (not the whole run), like an HTML span.

    python-pptx has no intra-run formatting, so each paragraph is rebuilt:
    its full text is split on the emphasis tokens, and each segment becomes its
    own run — plain runs keep the base style, token runs get the emphasis color.
    """
    if not emphasis:
        return
    tokens = [e.get("text", "") for e in emphasis if e.get("text")]
    for para in tf.paragraphs:
        full = "".join(r.text for r in para.runs)
        if not full or not any(tok in full for tok in tokens):
            continue
        # capture the paragraph's base style from its first run
        base = para.runs[0]
        base_size = base.font.size
        base_bold = base.font.bold
        base_color = base.font.color.rgb if base.font.color and base.font.color.type is not None else TEXT
        # clear existing runs
        for r in list(para.runs):
            r._r.getparent().remove(r._r)
        for seg in _split_on_tokens(full, tokens):
            run = para.add_run()
            run.text = seg
            emp = _emphasis_lookup(seg, emphasis)
            if emp is None:
                _set_run(run, size=base_size or BODY_PT, bold=bool(base_bold), color=base_color)
            else:
                kind = emp.get("kind", "main")
                if kind in EMPHASIS:
                    _set_run(run, size=base_size or BODY_PT, bold=True, color=EMPHASIS[kind])
                elif kind == "underline":
                    _set_run(run, size=base_size or BODY_PT, bold=bool(base_bold),
                             color=base_color, underline=True)
                elif kind == "invert":  # no run-level inversion; bold proxy
                    _set_run(run, size=base_size or BODY_PT, bold=True, color=base_color)
                else:
                    _set_run(run, size=base_size or BODY_PT, bold=bool(base_bold), color=base_color)


def _bottom_border_only(cell, width_pt=1.0, color="1A1A1A"):
    """Draw ONLY a bottom border on a table cell (no vertical/other lines).

    python-pptx has no high-level border API; we edit the cell's <a:tcPr>:
    add <a:lnB> (bottom) with a solid fill, and explicitly set the other three
    sides to noFill so the table style's grid lines don't show through.
    """
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
    # three sides: no line
    for tag in ("a:lnL", "a:lnR", "a:lnT"):
        ln = tcPr.makeelement(qn(tag), {"w": "12700"})
        ln.append(ln.makeelement(qn("a:noFill"), {}))
        tcPr.append(ln)
    # bottom: solid line
    w_emu = str(int(width_pt * 12700))  # 1pt = 12700 EMU
    lnB = tcPr.makeelement(qn("a:lnB"), {"w": w_emu, "cap": "flat"})
    fill = lnB.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
    fill.append(clr)
    lnB.append(fill)
    tcPr.append(lnB)


def _clear_table_style(table):
    """Set the table to the 'No Style, No Grid' built-in so default grid lines
    don't fight our per-cell bottom borders."""
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = tbl.makeelement(qn("a:tblPr"), {})
        tbl.insert(0, tblPr)
    # remove banding flags
    for attr in ("firstRow", "bandRow"):
        if attr in tblPr.attrib:
            del tblPr.attrib[attr]
    styleId = tblPr.find(qn("a:tableStyleId"))
    if styleId is None:
        styleId = tblPr.makeelement(qn("a:tableStyleId"), {})
        tblPr.append(styleId)
    styleId.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid


def _add_table(slide, spec):
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    n_rows = len(rows) + (1 if headers else 0)
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    if n_rows == 0 or n_cols == 0:
        return
    gf = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(spec["x_in"]), Inches(spec["y_in"]),
        Inches(spec["w_in"]), Inches(spec["h_in"]),
    )
    table = gf.table
    _clear_table_style(table)
    emphasize = spec.get("emphasize_row")  # 1-based among data rows

    r = 0
    if headers:
        for c, text in enumerate(headers):
            cell = table.cell(r, c)
            cell.fill.background()  # transparent (base color shows)
            cell.text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(text)
            _set_run(run, size=BODY_PT, bold=True, color=TEXT)
            _bottom_border_only(cell, width_pt=2.0)  # thead: 2px
        r += 1

    for di, row in enumerate(rows):
        is_emp = emphasize is not None and (di + 1) == emphasize
        last = di == len(rows) - 1
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.fill.background()
            cell.text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(row[c]) if c < len(row) else ""
            _set_run(run, size=BODY_PT, bold=is_emp,
                     color=MAIN if is_emp else TEXT)
            _bottom_border_only(cell, width_pt=2.0 if last else 1.0)
        r += 1


def _build_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BASE

    emphasis = spec.get("emphasis")
    # Content width in inches, used to estimate how many lines a string wraps to.
    content_w_in = 13.333 - 2 * 0.5

    # Title
    title_tf = _add_textbox(slide, MARGIN, MARGIN, CONTENT_W, Inches(1.2))
    run = title_tf.paragraphs[0].add_run()
    run.text = spec["title"]
    _set_run(run, size=TITLE_PT, bold=True, color=TEXT)
    _apply_emphasis(title_tf, emphasis)
    y_in = 1.5  # cursor in inches, just below the title

    # Subtitle (emphasis applies here too; advance the cursor by its wrapped height)
    if spec.get("subtitle"):
        sub = str(spec["subtitle"])
        sub_lines = _estimate_lines(sub, HEAD_PT, content_w_in)
        sub_h = 0.55 * sub_lines + 0.15
        sub_tf = _add_textbox(slide, MARGIN, Inches(y_in), CONTENT_W, Inches(sub_h))
        run = sub_tf.paragraphs[0].add_run()
        run.text = sub
        _set_run(run, size=HEAD_PT, bold=False, color=TEXT)
        _apply_emphasis(sub_tf, emphasis)
        y_in += sub_h + 0.25

    # Bullets (start below whatever the title+subtitle consumed; overridable)
    if spec.get("bullets"):
        by = spec.get("bullets_y_in", y_in)
        body_tf = _add_textbox(slide, MARGIN, Inches(by), CONTENT_W, Inches(7.5 - by - 0.6))
        for i, line in enumerate(spec["bullets"]):
            para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            para.space_after = Pt(14)
            run = para.add_run()
            run.text = "・" + str(line)
            _set_run(run, size=BODY_PT, color=TEXT)
        _apply_emphasis(body_tf, emphasis)
        y_in = by + 0.55 * len(spec["bullets"]) + 0.3

    # Body paragraph
    if spec.get("body"):
        by = spec.get("body_y_in", y_in)
        body_tf = _add_textbox(slide, MARGIN, Inches(by), CONTENT_W, Inches(7.5 - by - 0.6))
        run = body_tf.paragraphs[0].add_run()
        run.text = spec["body"]
        _set_run(run, size=BODY_PT, color=TEXT)
        _apply_emphasis(body_tf, emphasis)

    # Table
    if spec.get("table"):
        _add_table(slide, spec["table"])

    # Images (GPT-Image / theSVG parts) — placed directly, never via HTML
    for img in spec.get("images", []):
        path = img["_resolved"]
        kwargs = {"width": Inches(img["w_in"])}
        if img.get("h_in"):
            kwargs["height"] = Inches(img["h_in"])
        slide.shapes.add_picture(str(path), Inches(img["x_in"]),
                                 Inches(img["y_in"]), **kwargs)

    # Page number (bottom-right)
    if spec.get("page_number") is not None:
        pn_tf = _add_textbox(slide, SLIDE_W - Inches(1.2),
                             SLIDE_H - Inches(0.6), Inches(0.9), Inches(0.4))
        pn_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        run = pn_tf.paragraphs[0].add_run()
        run.text = str(spec["page_number"])
        _set_run(run, size=NUM_PT, color=TEXT)


def main():
    ap = argparse.ArgumentParser(description="Build a native .pptx from a JSON spec.")
    ap.add_argument("--spec", required=True, help="JSON spec path, or - for stdin")
    ap.add_argument("--out", required=True, help="Output .pptx path")
    args = ap.parse_args()

    if args.spec == "-":
        deck = json.load(sys.stdin)
        spec_dir = Path.cwd()
    else:
        spec_path = Path(args.spec)
        deck = json.loads(spec_path.read_text(encoding="utf-8"))
        spec_dir = spec_path.resolve().parent

    slides = deck.get("slides", [])
    if not slides:
        print("[build_pptx] ERROR: spec has no slides.", file=sys.stderr)
        sys.exit(1)

    # Resolve + validate image paths up front (fail fast with a clear message).
    for si, spec in enumerate(slides, 1):
        for img in spec.get("images", []):
            p = (spec_dir / img["path"]).resolve()
            if not p.is_file():
                print(f"[build_pptx] ERROR: slide {si} image not found: {p}",
                      file=sys.stderr)
                sys.exit(1)
            img["_resolved"] = p

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for spec in slides:
        if "title" not in spec:
            print("[build_pptx] ERROR: every slide needs a 'title'.", file=sys.stderr)
            sys.exit(1)
        _build_slide(prs, spec)

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"  Built: {out}  ({len(slides)} slide(s))", file=sys.stderr)


if __name__ == "__main__":
    main()
